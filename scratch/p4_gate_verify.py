"""P4-GATE — FULL pipeline verification (Postgres + Qdrant).

AGENTS.md §7: ALWAYS run through DocumentRepository.upsert_document()
→ ContentBlock rows → tsvector → lexical search works.
"""
from __future__ import annotations

import asyncio, io, time, uuid
from apps.api.config import settings
from rekanvault.contracts.documents import DocumentBlock, DocumentLocator, DocumentVersion, NormalizedDocument
from rekanvault.contracts.documents import SourceProvider
from rekanvault.evidence.embedding import EmbeddingService
from rekanvault.evidence.indexing import IndexingPipeline
from rekanvault.evidence.chunker import Chunker
from rekanvault.evidence.retrieval import RetrievalPipeline
from rekanvault.evaluation.runner import EvaluationRunner, load_golden_questions, _count_targets, _idcg
from rekanvault.storage.database import init_db, get_db_session
from rekanvault.storage.document_repo import DocumentRepository
from rekanvault.storage.qdrant import QdrantStore
import google.auth.transport.requests, google.oauth2.credentials, googleapiclient.discovery


def _build_drive_service():
    creds = google.oauth2.credentials.Credentials(
        token=None, refresh_token=settings.RV_GOOGLE_PILOT_REFRESH_TOKEN,
        token_uri="https://oauth2.googleapis.com/token",
        client_id=settings.RV_GOOGLE_CLIENT_ID, client_secret=settings.RV_GOOGLE_CLIENT_SECRET,
        scopes=settings.RV_GOOGLE_OAUTH_SCOPES.split(","),
    )
    creds.refresh(google.auth.transport.requests.Request())
    return googleapiclient.discovery.build("drive", "v3", credentials=creds)


def _list_all_files(service, folder_id, prefix=""):
    results = []
    pt = None
    while True:
        resp = service.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            fields="nextPageToken, files(id, name, mimeType, size)", pageSize=100, pageToken=pt,
        ).execute()
        for f in resp.get("files", []):
            if f["mimeType"] == "application/vnd.google-apps.folder":
                results.extend(_list_all_files(service, f["id"], f"{prefix}{f['name']}/"))
            elif int(f.get("size", 0)) <= settings.RV_MAX_SOURCE_FILE_BYTES:
                f["_path"] = f"{prefix}{f['name']}"
                results.append(f)
        pt = resp.get("nextPageToken")
        if not pt: break
    return results


def _download_text(service, file_id, mime_type):
    try:
        if mime_type == "application/vnd.google-apps.document":
            resp = service.files().export(fileId=file_id, mimeType="text/plain").execute()
            return resp.decode("utf-8") if isinstance(resp, bytes) else str(resp)
        elif mime_type == "application/vnd.google-apps.presentation":
            resp = service.files().export(fileId=file_id, mimeType="text/plain").execute()
            return resp.decode("utf-8") if isinstance(resp, bytes) else str(resp)
        elif mime_type == "application/vnd.google-apps.spreadsheet":
            resp = service.files().export(fileId=file_id, mimeType="text/csv").execute()
            return resp.decode("utf-8") if isinstance(resp, bytes) else str(resp)
        elif "presentation" in mime_type:
            resp = service.files().get_media(fileId=file_id).execute()
            from pptx import Presentation
            prs = Presentation(io.BytesIO(resp))
            return "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes
                if shape.has_text_frame for p in shape.text_frame.paragraphs if p.text.strip()
            )
        elif "pdf" in mime_type:
            resp = service.files().get_media(fileId=file_id).execute()
            import pdfplumber
            parts = []
            with pdfplumber.open(io.BytesIO(resp)) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t: parts.append(t)
            return "\n\n".join(parts) if parts else None
        elif "document" in mime_type:
            resp = service.files().get_media(fileId=file_id).execute()
            from docx import Document as DocxDocument
            doc = DocxDocument(io.BytesIO(resp))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif "spreadsheet" in mime_type:
            resp = service.files().export(fileId=file_id, mimeType="text/csv").execute()
            return resp.decode("utf-8") if isinstance(resp, bytes) else str(resp)
        else:
            resp = service.files().get_media(fileId=file_id).execute()
            try:
                return resp.decode("utf-8")
            except (UnicodeDecodeError, AttributeError):
                return None
    except Exception as e:
        print(f"    download err: {e}"); return None


async def main():
    t_start = time.time()
    init_db()
    embed = EmbeddingService()
    qdrant = QdrantStore(settings)
    # Drop old data
    if await qdrant.client.collection_exists(settings.RV_QDRANT_COLLECTION):
        await qdrant.client.delete_collection(settings.RV_QDRANT_COLLECTION)
    await qdrant.ensure_collection()
    doc_repo = DocumentRepository()
    chunker = Chunker(repo=doc_repo)
    ws_id = uuid.UUID(settings.RV_PILOT_WORKSPACE_ID)
    src_id = uuid.UUID(settings.RV_PILOT_SOURCE_ID)

    # ── Stage 1: Sync ──
    print("=== STAGE 1: Google Drive Sync ===")
    service = _build_drive_service()
    files = _list_all_files(service, settings.RV_GOOGLE_FOLDER_ID)
    print(f"Found {len(files)} files")

    # ── Stage 2: Postgres + Qdrant Ingestion ──
    print(f"\n=== STAGE 2: Postgres → Chunk → Embed → Qdrant ===")
    total_chunks = 0
    async for session in get_db_session():
        for i, f in enumerate(files):
            name, doc_id = f["name"], f["id"]
            path = f.get("_path", name)
            text = _download_text(service, doc_id, f["mimeType"])
            if not text:
                continue

            # Build NormalizedDocument
            block = DocumentBlock(block_id=f"{doc_id}-b0", block_type="text", content=text, sequence=0)
            version = DocumentVersion(
                version_id=f"{doc_id}-v1", document_id=doc_id, version_number=1,
                content_hash=f"{doc_id}-hash", blocks=[block],
            )
            norm = NormalizedDocument(
                document_id=doc_id, workspace_id=str(ws_id), source_id=str(src_id),
                title=path, provider=SourceProvider.GOOGLE_DRIVE,
                locator=DocumentLocator(provider=SourceProvider.GOOGLE_DRIVE, native_id=doc_id, mime_type=f.get("mimeType"), uri=f"https://drive.google.com/file/d/{doc_id}/view"),
                active_version_id=f"{doc_id}-v1", versions=[version],
                metadata={"storage_path": None},
            )

            # Write to Postgres → ContentBlock rows → tsvector populated
            await doc_repo.upsert_document(session, ws_id, src_id, norm)

            # Get the freshly inserted document + version for indexing
            doc = await doc_repo.get_by_external_id(session, ws_id, src_id, doc_id)
            if doc is None:
                continue
            dv = await doc_repo.get_latest_version(session, doc.id)
            if dv is None:
                continue

            # Chunk → embed → Qdrant
            pipeline = IndexingPipeline(session=session, chunker=chunker, embed=embed, qdrant=qdrant, doc_repo=doc_repo)
            n = await pipeline.index_version(dv.id)
            total_chunks += n
            print(f"  [{i+1}/{len(files)}] {path[:60]}... → {n} chunks")

        await session.commit()
        break  # one session

    sync_elapsed = time.time() - t_start
    print(f"\nIndexed {total_chunks} chunks in {sync_elapsed:.0f}s")

    if total_chunks == 0:
        print("No chunks — aborting eval.")
        await qdrant.close(); return

    # ── Stage 3: Verify lexical has data ──
    async for session in get_db_session():
        from sqlalchemy import text
        r = await session.execute(text("SELECT count(*) FROM content_blocks"))
        cb_count = r.scalar()
        print(f"ContentBlocks in Postgres: {cb_count}")
        await session.commit(); break

    # ── Stage 4: Evaluate (full hybrid pipeline, no negatives) ──
    print(f"\n=== STAGE 4: Golden Set Evaluation ===")
    all_qs = load_golden_questions("docs/REKANVAULT_GOLDEN_SET.md")
    qs = [q for q in all_qs if q["category"] not in ("NEGATIVE", "INSUFFICIENT")]
    print(f"Running {len(qs)}/{len(all_qs)} questions (full hybrid, 1s delay, negatives skipped)...")
    async for session in get_db_session():
        pipeline = RetrievalPipeline(session=session, embed=embed, qdrant=qdrant)
        runner = EvaluationRunner(pipeline=pipeline)
        import math as _math
        hits_total, rr_total, dcg_total = 0, 0.0, 0.0
        cats: dict[str, dict[str, int]] = {}
        t0 = time.time()
        for i, qq in enumerate(qs):
            category = qq["category"]
            try:
                detail = await runner.evaluate_question(qq, top_k=10)
                ranks = detail["ranks"]
                correct = detail["correct"]
                first_rank = detail["first_rank"]
                n_targets = detail["n_targets"]
            except Exception:
                ranks, correct, first_rank, n_targets = [], None, None, 0
            cats.setdefault(category, {"t": 0, "h": 0})
            cats[category]["t"] += 1
            if correct is True:
                hits_total += 1
                cats[category]["h"] += 1
            if first_rank is not None:
                rr_total += 1.0 / first_rank
            for r in ranks:
                dcg_total += 1.0 / _math.log2(r + 1)
            if (i + 1) % 30 == 0:
                print(f"  {i+1}/{len(qs)}...")
            await asyncio.sleep(1.0)
        elapsed = time.time() - t0
        await session.commit(); break

    n = len(qs)
    total_elapsed = time.time() - t_start
    recall = hits_total / n if n else 0.0
    mrr = rr_total / n if n else 0.0
    idcg_sum = sum(_idcg(_count_targets(qq.get("target_source", ""))) for qq in qs)
    ndcg = dcg_total / idcg_sum if idcg_sum > 0 else (1.0 if dcg_total == 0 else 0.0)
    print(f"\n{'='*60}")
    print(f"FULL-PIPELINE P4-GATE — {n} scorable questions (negatives excluded)")
    print(f"{'='*60}")
    print(f"ContentBlocks:  {cb_count}")
    print(f"Chunks indexed: {total_chunks}")
    print(f"Recall@10:      {recall:.4f}  (target >= 0.85)")
    print(f"MRR:            {mrr:.4f}")
    print(f"nDCG@10:        {ndcg:.4f}")
    print(f"Hits:           {hits_total}/{n} ({hits_total/n*100:.1f}%)" if n else "Hits: N/A")
    print(f"Ingest time:    {sync_elapsed:.0f}s")
    print(f"Eval time:      {elapsed:.0f}s")
    print(f"Total:          {total_elapsed:.0f}s")
    print(f"\nBy category:")
    for cat, c in sorted(cats.items()):
        pct = c["h"] / c["t"] * 100 if c["t"] else 0
        mark = "✓" if pct >= 90 else "~" if pct >= 50 else "✗"
        print(f"  {cat:15s}: {c['h']:>3}/{c['t']:<3} ({pct:.0f}%) {mark}")

    await qdrant.close()


if __name__ == "__main__":
    asyncio.run(main())
